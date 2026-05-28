"""
generate-pptx.py
-----------------
Build "Pitch Deck.pptx" from "Pitch Deck.md".

Uses python-pptx. Same markdown source as generate-pdf.py — same 14 slides,
same accent palette, same content. PowerPoint output is intentionally
simpler than the HTML/PDF rendering: clean type, real PPTX tables, accent
bar at the top of each slide.

Re-run any time Pitch Deck.md changes.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT      = Path(__file__).resolve().parent
VAULT_DIR = ROOT.parent
DECK_MD   = VAULT_DIR / "Pitch Deck.md"
PPTX_OUT  = ROOT / "Pitch Deck.pptx"

# ---------------------------------------------------------------------------
# Design tokens — ported from styles/deck.css
# ---------------------------------------------------------------------------
COLOR_BG_DARK       = RGBColor(0x1A, 0x16, 0x12)
COLOR_BG_LIGHT      = RGBColor(0xFA, 0xF7, 0xF1)
COLOR_SURFACE_DARK  = RGBColor(0x2A, 0x24, 0x1F)
COLOR_SURFACE_LIGHT = RGBColor(0xF0, 0xEA, 0xDF)
COLOR_TEXT_DARK     = RGBColor(0xF0, 0xE9, 0xDD)
COLOR_TEXT_LIGHT    = RGBColor(0x22, 0x1E, 0x1A)
COLOR_MUTED_DARK    = RGBColor(0xA8, 0x9F, 0x90)
COLOR_MUTED_LIGHT   = RGBColor(0x66, 0x60, 0x58)
COLOR_ACCENT_EMBER  = RGBColor(0xE5, 0x62, 0x1A)
COLOR_ACCENT_TEAL   = RGBColor(0x2D, 0x8E, 0x9E)
COLOR_ACCENT_INDIGO = RGBColor(0x53, 0x4A, 0xB0)
COLOR_ACCENT_AMBER  = RGBColor(0xF5, 0xB2, 0x3A)

FONT_DISPLAY = "Inter"        # Space Grotesk if installed locally
FONT_BODY    = "Inter"
FONT_MONO    = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Markdown parsing — same shape as generate-pdf.py
# ---------------------------------------------------------------------------
SLIDE_HEADER_RE = re.compile(
    r"^---\s*\nslide:\s*(\d+)\s*\n([\s\S]*?)^---\s*\n",
    re.MULTILINE,
)


def parse_frontmatter(block: str) -> dict[str, str]:
    out: dict[str, str] = {}
    for line in block.strip().splitlines():
        if ":" in line:
            k, _, v = line.partition(":")
            out[k.strip()] = v.strip().strip('"').strip("'")
    return out


def split_slides(text: str) -> list[tuple[dict[str, str], str]]:
    slides: list[tuple[dict[str, str], str]] = []
    matches = list(SLIDE_HEADER_RE.finditer(text))
    for i, m in enumerate(matches):
        fm = parse_frontmatter(m.group(2))
        body_start = m.end()
        body_end   = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[body_start:body_end].strip()
        body = re.sub(r"^# Appendix.*$", "", body, flags=re.MULTILINE | re.DOTALL)
        slides.append((fm, body))
    return slides


# ---------------------------------------------------------------------------
# Markdown subsection parser — pulls (heading, paragraphs, table, code, lists)
# out of a slide body in document order
# ---------------------------------------------------------------------------
CARD_OPEN_RE  = re.compile(r"^:::\s*card(.*)$")
CARD_CLOSE_RE = re.compile(r"^:::\s*$")


def parse_attrs(line: str) -> dict[str, str]:
    return dict(re.findall(r'(\w+)="([^"]*)"', line))


def parse_blocks(body: str) -> list[dict]:
    """Return a list of typed blocks: heading | paragraph | bullets | table | code | quote | card-group."""
    blocks: list[dict] = []
    lines = body.splitlines()
    i = 0
    while i < len(lines):
        ln = lines[i]
        s  = ln.strip()
        if not s:
            i += 1
            continue
        # heading
        if s.startswith("# "):
            blocks.append({"type": "h1", "text": s[2:].strip()})
            i += 1; continue
        if s.startswith("## "):
            blocks.append({"type": "h2", "text": s[3:].strip()})
            i += 1; continue
        # blockquote
        if s.startswith("> "):
            quote = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote.append(lines[i].strip().lstrip(">").strip())
                i += 1
            blocks.append({"type": "quote", "lines": quote})
            continue
        # code fence
        if s.startswith("```"):
            i += 1
            code = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            blocks.append({"type": "code", "lines": code})
            continue
        # table — line with '|' followed by separator row
        if "|" in s and i + 1 < len(lines) and re.match(r"^\s*\|?\s*:?-+", lines[i + 1].strip()):
            table: list[list[str]] = []
            while i < len(lines) and "|" in lines[i]:
                row = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not re.match(r"^:?-+:?$", row[0]):
                    table.append(row)
                i += 1
            blocks.append({"type": "table", "rows": table})
            continue
        # bullets
        if s.startswith("- ") or s.startswith("* "):
            bullets = []
            while i < len(lines) and (lines[i].strip().startswith("- ") or lines[i].strip().startswith("* ")):
                bullets.append(lines[i].strip()[2:])
                i += 1
            blocks.append({"type": "bullets", "items": bullets})
            continue
        # numbered list
        if re.match(r"^\d+\.\s", s):
            items = []
            while i < len(lines) and re.match(r"^\d+\.\s", lines[i].strip()):
                items.append(re.sub(r"^\d+\.\s", "", lines[i].strip()))
                i += 1
            blocks.append({"type": "ordered", "items": items})
            continue
        # card group
        if CARD_OPEN_RE.match(s):
            cards = []
            while i < len(lines) and CARD_OPEN_RE.match(lines[i].strip()):
                attrs = parse_attrs(CARD_OPEN_RE.match(lines[i].strip()).group(1))
                i += 1
                card_body = []
                while i < len(lines) and not CARD_CLOSE_RE.match(lines[i].strip()):
                    card_body.append(lines[i])
                    i += 1
                i += 1  # consume :::
                cards.append({"meta": attrs, "body": parse_blocks("\n".join(card_body))})
                while i < len(lines) and lines[i].strip() == "":
                    i += 1
            blocks.append({"type": "card-group", "cards": cards})
            continue
        # paragraph
        para = [ln]
        i += 1
        while i < len(lines) and lines[i].strip() and not _starts_new_block(lines[i].strip()):
            para.append(lines[i])
            i += 1
        blocks.append({"type": "p", "text": " ".join(p.strip() for p in para)})

    return blocks


def _starts_new_block(s: str) -> bool:
    return (
        s.startswith("#")
        or s.startswith("> ")
        or s.startswith("```")
        or s.startswith("- ")
        or s.startswith("* ")
        or s.startswith(":::")
        or bool(re.match(r"^\d+\.\s", s))
        or ("|" in s)
    )


# ---------------------------------------------------------------------------
# Inline markdown — only need bold runs
# ---------------------------------------------------------------------------
INLINE_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
WIKILINK_RE    = re.compile(r"\[\[([^|\]]+)(?:\|([^\]]+))?\]\]")


def emit_runs(paragraph, text: str, *, color: RGBColor, font: str, size: int, base_bold: bool = False) -> None:
    """Split text by ** and emit runs preserving bold."""
    text = WIKILINK_RE.sub(lambda m: m.group(2) or m.group(1), text)
    pos = 0
    for m in INLINE_BOLD_RE.finditer(text):
        if m.start() > pos:
            _run(paragraph, text[pos:m.start()], color=color, font=font, size=size, bold=base_bold)
        _run(paragraph, m.group(1), color=color, font=font, size=size, bold=True)
        pos = m.end()
    if pos < len(text):
        _run(paragraph, text[pos:], color=color, font=font, size=size, bold=base_bold)


def _run(paragraph, text: str, *, color: RGBColor, font: str, size: int, bold: bool = False, italic: bool = False) -> None:
    r = paragraph.add_run()
    r.text = text
    r.font.name  = font
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def add_blank_slide(prs: Presentation, theme: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK if theme == "dark" else COLOR_BG_LIGHT
    bg.shadow.inherit = False
    return slide


def add_accent_bar(slide, color: RGBColor) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color


def add_footer(slide, slide_no: int, total: int, theme: str) -> None:
    muted = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT
    # Brand mark left
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(5), Inches(0.3))
    tf = tb.text_frame; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    _run(p, "GALVON · MANUFACTURING AI", color=muted, font=FONT_MONO, size=8)
    # Slide number right
    tb2 = slide.shapes.add_textbox(Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3))
    tf2 = tb2.text_frame; tf2.margin_top = 0; tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    _run(p2, f"{slide_no:02d} / {total:02d}", color=muted, font=FONT_MONO, size=8)


def add_textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0
    return tb, tf


def render_blocks_into_textbox(
    tf,
    blocks: list[dict],
    theme: str,
    *,
    body_size: int = 14,
    bullet_size: int = 14,
) -> None:
    text     = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    muted    = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT
    accent   = COLOR_ACCENT_AMBER if theme == "dark" else COLOR_ACCENT_EMBER
    first = True
    for b in blocks:
        if b["type"] == "h2":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_after = Pt(8)
            emit_runs(p, b["text"], color=accent, font=FONT_DISPLAY, size=22, base_bold=True)
        elif b["type"] == "p":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_after = Pt(6)
            emit_runs(p, b["text"], color=text, font=FONT_BODY, size=body_size)
        elif b["type"] == "bullets":
            for item in b["items"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.space_after = Pt(4)
                _run(p, "•  ", color=accent, font=FONT_BODY, size=bullet_size, bold=True)
                emit_runs(p, item, color=text, font=FONT_BODY, size=bullet_size)
                first = False
            continue
        elif b["type"] == "ordered":
            for n, item in enumerate(b["items"], 1):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.space_after = Pt(4)
                _run(p, f"{n}.  ", color=accent, font=FONT_DISPLAY, size=bullet_size, bold=True)
                emit_runs(p, item, color=text, font=FONT_BODY, size=bullet_size)
                first = False
            continue
        elif b["type"] == "quote":
            for line in b["lines"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.space_after = Pt(4)
                _run(p, line, color=muted, font=FONT_BODY, size=body_size, italic=True)
                first = False
            continue
        elif b["type"] == "code":
            for line in b["lines"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.space_after = Pt(0)
                _run(p, line if line else " ", color=muted, font=FONT_MONO, size=10)
                first = False
            continue
        first = False


def add_table(slide, rows: list[list[str]], x, y, w, h, theme: str) -> None:
    n_rows = len(rows)
    n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    text   = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    muted  = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT
    accent = COLOR_ACCENT_AMBER if theme == "dark" else COLOR_ACCENT_EMBER
    bg     = COLOR_BG_DARK if theme == "dark" else COLOR_BG_LIGHT

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            p = tf.paragraphs[0]
            if ri == 0:
                _run(p, val, color=muted, font=FONT_DISPLAY, size=10, bold=True)
            elif ci == len(row) - 1 and len(row) > 1:
                emit_runs(p, val, color=accent, font=FONT_MONO, size=11, base_bold=True)
            else:
                emit_runs(p, val, color=text, font=FONT_BODY, size=11)


# ---------------------------------------------------------------------------
# Layout handlers
# ---------------------------------------------------------------------------
def build_layout_title(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    h2 = next((b for b in blocks if b["type"] == "h2"), None)
    para = next((b for b in blocks if b["type"] == "p"), None)
    quote = next((b for b in blocks if b["type"] == "quote"), None)

    text  = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    muted = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT

    tb, tf = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(12), Inches(3))
    p1 = tf.paragraphs[0]
    if h1:
        _run(p1, h1["text"], color=text, font=FONT_DISPLAY, size=72, bold=True)
    if h2:
        p2 = tf.add_paragraph(); p2.space_before = Pt(6)
        _run(p2, h2["text"], color=muted, font=FONT_DISPLAY, size=40)

    if para:
        tbp, tfp = add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.5))
        emit_runs(tfp.paragraphs[0], para["text"], color=text, font=FONT_BODY, size=16)
    if quote:
        tbq, tfq = add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5))
        for line in quote["lines"]:
            _run(tfq.paragraphs[0], line, color=muted, font=FONT_MONO, size=11)


def build_layout_full_bleed(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    tb, tf = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(1.5))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=56, bold=True)
    # body paragraphs
    body_blocks = [b for b in blocks if b["type"] not in ("h1",)]
    tb2, tf2 = add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(4.0))
    render_blocks_into_textbox(tf2, body_blocks, theme, body_size=18)


def build_layout_stats(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    h2 = next((b for b in blocks if b["type"] == "h2"), None)
    para_blocks = [b for b in blocks if b["type"] == "p"]
    table = next((b for b in blocks if b["type"] == "table"), None)

    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    muted = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT

    tb, tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=36, bold=True)

    if para_blocks:
        tbp, tfp = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8))
        emit_runs(tfp.paragraphs[0], para_blocks[0]["text"], color=muted, font=FONT_BODY, size=16)

    if table:
        add_table(slide, table["rows"], Inches(0.8), Inches(2.5), Inches(11.5), Inches(4.0), theme)

    # trailing paragraphs after the table
    if len(para_blocks) > 1:
        tb3, tf3 = add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6))
        emit_runs(tf3.paragraphs[0], para_blocks[-1]["text"], color=muted, font=FONT_BODY, size=12, base_bold=False)


def build_layout_split(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT

    tb, tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=36, bold=True)

    rest = [b for b in blocks if b["type"] != "h1"]
    code_blocks = [b for b in rest if b["type"] == "code"]
    other_blocks = [b for b in rest if b["type"] != "code"]

    if code_blocks:
        tb2, tf2 = add_textbox(slide, Inches(0.8), Inches(1.7), Inches(7.5), Inches(5.0))
        render_blocks_into_textbox(tf2, code_blocks, theme)
        tb3, tf3 = add_textbox(slide, Inches(8.6), Inches(1.7), Inches(4.0), Inches(5.0))
        render_blocks_into_textbox(tf3, other_blocks, theme, body_size=13)
    else:
        tb2, tf2 = add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0))
        render_blocks_into_textbox(tf2, rest, theme, body_size=14)


def build_layout_cards(slide, blocks, theme, n_cards: int):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    p_pre = next((b for b in blocks if b["type"] == "p"), None)
    cards_block = next((b for b in blocks if b["type"] == "card-group"), None)
    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    muted = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT

    tb, tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=34, bold=True)

    if p_pre:
        tbp, tfp = add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8))
        emit_runs(tfp.paragraphs[0], p_pre["text"], color=muted, font=FONT_BODY, size=14)

    if not cards_block:
        return

    cards = cards_block["cards"][:n_cards]
    gap = Inches(0.25)
    total_w = Inches(11.5)
    card_w = Emu(int((total_w - gap * (len(cards) - 1)) / max(len(cards), 1)))
    card_h = Inches(4.6)
    x = Inches(0.8)
    y = Inches(2.3)

    accent_map = {
        "ember": COLOR_ACCENT_EMBER,
        "teal":  COLOR_ACCENT_TEAL,
        "indigo": COLOR_ACCENT_INDIGO,
    }

    for card in cards:
        meta = card["meta"]
        accent = accent_map.get(meta.get("accent", "ember"), COLOR_ACCENT_EMBER)
        # card surface
        surface = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        surface.line.fill.background()
        surface.fill.solid()
        surface.fill.fore_color.rgb = COLOR_SURFACE_DARK if theme == "dark" else COLOR_SURFACE_LIGHT
        surface.adjustments[0] = 0.04
        surface.shadow.inherit = False
        # accent bar on top
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.08))
        accent_bar.line.fill.background()
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        # title + body
        tb_c, tf_c = add_textbox(slide, x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.5))
        title = meta.get("title", "")
        if title:
            p_t = tf_c.paragraphs[0]
            _run(p_t, title, color=text, font=FONT_DISPLAY, size=18, bold=True)
            p_spacer = tf_c.add_paragraph(); p_spacer.space_after = Pt(6); _run(p_spacer, " ", color=text, font=FONT_BODY, size=6)
        render_blocks_into_textbox_existing(tf_c, card["body"], theme, body_size=12, bullet_size=12)
        x = x + card_w + gap


def render_blocks_into_textbox_existing(tf, blocks, theme, body_size=12, bullet_size=12):
    """Like render_blocks_into_textbox but assumes the first paragraph is already used."""
    text  = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    muted = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT
    accent = COLOR_ACCENT_AMBER if theme == "dark" else COLOR_ACCENT_EMBER
    for b in blocks:
        if b["type"] == "p":
            p = tf.add_paragraph(); p.space_after = Pt(4)
            emit_runs(p, b["text"], color=text, font=FONT_BODY, size=body_size)
        elif b["type"] == "bullets":
            for item in b["items"]:
                p = tf.add_paragraph(); p.space_after = Pt(2)
                _run(p, "•  ", color=accent, font=FONT_BODY, size=bullet_size, bold=True)
                emit_runs(p, item, color=text, font=FONT_BODY, size=bullet_size)
        elif b["type"] == "h2":
            p = tf.add_paragraph(); p.space_after = Pt(4)
            emit_runs(p, b["text"], color=accent, font=FONT_DISPLAY, size=14, base_bold=True)


def build_layout_diagram(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT

    tb, tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=34, bold=True)

    rest = [b for b in blocks if b["type"] != "h1"]
    tb2, tf2 = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.3))
    render_blocks_into_textbox(tf2, rest, theme)


def build_layout_story(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    tb, tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=32, bold=True)
    rest = [b for b in blocks if b["type"] != "h1"]
    tb2, tf2 = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.3))
    render_blocks_into_textbox(tf2, rest, theme, body_size=13)


def build_layout_table(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    p_block = next((b for b in blocks if b["type"] == "p"), None)
    table = next((b for b in blocks if b["type"] == "table"), None)
    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    muted = COLOR_MUTED_DARK if theme == "dark" else COLOR_MUTED_LIGHT

    tb, tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=32, bold=True)
    y = Inches(1.5)
    if p_block:
        tbp, tfp = add_textbox(slide, Inches(0.8), y, Inches(11.5), Inches(0.7))
        emit_runs(tfp.paragraphs[0], p_block["text"], color=muted, font=FONT_BODY, size=13)
        y = Inches(2.2)
    if table:
        add_table(slide, table["rows"], Inches(0.8), y, Inches(11.5), Inches(4.5), theme)


def build_layout_cta(slide, blocks, theme):
    add_accent_bar(slide, COLOR_ACCENT_EMBER)
    h1 = next((b for b in blocks if b["type"] == "h1"), None)
    text = COLOR_TEXT_DARK if theme == "dark" else COLOR_TEXT_LIGHT
    tb, tf = add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(1.5))
    if h1:
        _run(tf.paragraphs[0], h1["text"], color=text, font=FONT_DISPLAY, size=52, bold=True)
    rest = [b for b in blocks if b["type"] != "h1"]
    tb2, tf2 = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(3.8))
    render_blocks_into_textbox(tf2, rest, theme, body_size=15, bullet_size=15)


LAYOUT_HANDLERS = {
    "title":        build_layout_title,
    "full-bleed":   build_layout_full_bleed,
    "stats":        build_layout_stats,
    "split":        build_layout_split,
    "two-cards":    lambda s, b, t: build_layout_cards(s, b, t, n_cards=2),
    "three-cards":  lambda s, b, t: build_layout_cards(s, b, t, n_cards=3),
    "diagram":      build_layout_diagram,
    "chart":        build_layout_diagram,
    "story":        build_layout_story,
    "tier-table":   build_layout_table,
    "technical":    build_layout_table,
    "cta":          build_layout_cta,
}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main() -> int:
    if not DECK_MD.exists():
        sys.exit(f"[!] Source not found: {DECK_MD}")

    md_text = DECK_MD.read_text(encoding="utf-8")
    slides = split_slides(md_text)
    if not slides:
        sys.exit("[!] No slides parsed.")

    print(f"[*] Parsed {len(slides)} slides")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(slides)
    for idx, (fm, body) in enumerate(slides):
        layout = fm.get("layout", "default")
        theme  = fm.get("theme", "light")
        slide_no = int(fm.get("slide", str(idx + 1)))
        blocks = parse_blocks(body)

        slide = add_blank_slide(prs, theme)
        handler = LAYOUT_HANDLERS.get(layout, build_layout_split)
        handler(slide, blocks, theme)
        add_footer(slide, slide_no, total, theme)
        print(f"    slide {slide_no} ({layout}, {theme}) — {len(blocks)} blocks")

    prs.save(PPTX_OUT)
    print(f"[+] Wrote {PPTX_OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
